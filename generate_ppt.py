from pptx import Presentation

def replace_text_in_shape(shape, old_text, new_text):
    """Safely replaces text in a shape while trying to preserve the template's font formatting."""
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)

def inject_content(slide, title_keyword, bullet_points):
    """Finds the body placeholder in the template and adds our detailed text."""
    title_shape = slide.shapes.title
    if title_shape and title_keyword in title_shape.text.upper():
        clean_title = title_shape.text.split("<<")[0].strip()
        title_shape.text = clean_title

    body_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title_shape:
            if "…" in shape.text or "Abstract" in shape.text or "Content" in shape.text:
                body_shape = shape
                break
                
    if not body_shape:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break

    if body_shape:
        tf = body_shape.text_frame
        tf.clear() 
        
        p = tf.paragraphs[0]
        p.text = bullet_points[0][0]
        p.level = bullet_points[0][1]
        
        for pt, lvl in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = pt
            p.level = lvl

def fill_university_template():
    template_filename = 'RTRP__KR24_II-II_2026_Project_Presentation_Template_1.0.pptx'
    
    try:
        prs = Presentation(template_filename)
    except Exception as e:
        print(f"ERROR: Could not find '{template_filename}'.")
        return

    # ==========================================
    # SLIDE 1 (Index 0): TITLE SLIDE
    # ==========================================
    title_slide = prs.slides[0]
    for shape in title_slide.shapes:
        replace_text_in_shape(shape, "dd-Mon-yyyy", "02-Apr-2026")
        replace_text_in_shape(shape, "(Name & Roll number)", "Raghuveer\nAkshith\nRohith\nYagna sai")
        replace_text_in_shape(shape, "... (all members)", "")
        replace_text_in_shape(shape, "…", "")
        # You will still need to manually type your Mentor's name on the first slide!

    # ==========================================
    # SLIDE 2 (Index 1): CONTENTS
    # ==========================================
    inject_content(prs.slides[1], "CONTENTS", [
        ("1. Introduction", 0),
        ("2. Requirement (Features & Specifications)", 0),
        ("3. Design (Architecture & Database)", 0),
        ("4. Development / Implementation", 0),
        ("5. Project Status", 0),
        ("6. Learnings So Far", 0)
    ])

    # ==========================================
    # SLIDE 3 (Index 2): INTRODUCTION
    # ==========================================
    inject_content(prs.slides[2], "INTRODUCTION", [
        ("End User Requirement & Expected Outcome:", 0),
        ("Users require a fast, frictionless way to generate contextual practice quizzes from their specific study materials (PDFs/Text) to reduce manual study effort.", 1),
        ("Abstract / Gist of Problem Statement:", 0),
        ("Traditional assessment creation is highly manual and time-consuming, leading to content latency and 'blank page' friction for self-learners.", 1),
        ("Objective(s) & Solution Impact:", 0),
        ("Automate entire content creation pipeline using the Gemini 2.5 Flash LLM.", 1),
        ("Drive user engagement and continuous learning through competitive gamification (Tracking XP, Streaks, and Tier Ranks).", 1)
    ])

    # ==========================================
    # SLIDE 4 (Index 3): REQUIREMENT
    # ==========================================
    inject_content(prs.slides[3], "REQUIREMENT", [
        ("Core Features & Use Cases:", 0),
        ("Dynamic Quiz Generation: Converts direct text, topics, or PDF uploads into interactive JSON-driven cards.", 1),
        ("Gamified HUD: Real-time tracking of player stats (Academy Student to Hokage).", 1),
        ("Functional Requirements:", 0),
        ("Secure user authentication with Werkzeug password hashing.", 1),
        ("Asynchronous Fetch API calls to prevent page reloads during generation.", 1),
        ("Persistent local session management via Flask cookies.", 1),
        ("Non-Functional Requirements:", 0),
        ("High performance (<15s generation time) and strict JSON formatting for UI stability.", 1)
    ])

    # ==========================================
    # SLIDE 5 (Index 4): DESIGN
    # ==========================================
    inject_content(prs.slides[4], "DESIGN", [
        ("Technology Stack & Architecture:", 0),
        ("Frontend: HTML5, CSS3 (Glass-morphism animations), Vanilla JS.", 1),
        ("Backend: Python, Flask, google-genai SDK, PyPDF2 parser.", 1),
        ("Database Design (SQLite):", 0),
        ("Users Table: id, username, password_hash, score, xp, streak.", 1),
        ("Quizzes Table: id, topic, difficulty, questions_json.", 1),
        ("[INSERT IMAGE: PASTE YOUR ARCHITECTURE DIAGRAM HERE]", 0),
        ("[INSERT IMAGE: PASTE YOUR MOCKUP SCREENS HERE]", 0)
    ])

    # ==========================================
    # SLIDE 6 (Index 5): DEVELOPMENT/IMPLEMENTATION
    # ==========================================
    inject_content(prs.slides[5], "DEVELOPMENT", [
        ("Coding Considerations & Tools Used:", 0),
        ("Modular Python backend developed in VS Code using virtual environments (venv).", 1),
        ("Prompt Engineering utilized to enforce strict JSON array output from the Gemini API, preventing hallucinated formatting.", 1),
        ("User Roles & Permissions:", 0),
        ("Guest Mode: Can generate quizzes but stats reset upon leaving.", 1),
        ("Authenticated User: Persistent stats, rank saving, and secure login.", 1),
        ("Exceptions Handled:", 0),
        ("File parsing errors (non-text PDFs), API quota limits, and invalid login handling.", 1)
    ])

    # ==========================================
    # SLIDE 7 (Index 6): PROJECT STATUS
    # ==========================================
    inject_content(prs.slides[6], "STATUS", [
        ("Phases Completed:", 0),
        ("UI/UX layout and Gamification logic (Streak multipliers & XP).", 1),
        ("Backend API routing, PyPDF2 parsing, and Gemini Model integration.", 1),
        ("Database initialization, user accounts, and session persistence.", 1),
        ("Remaining Phases:", 0),
        ("Cloud Deployment to a PaaS (Platform as a Service) like Render.", 1),
        ("Implementation of true model fine-tuning based on accumulated user data.", 1),
        ("Target Completion Time: May 2026", 0)
    ])

    # ==========================================
    # SLIDE 8 (Index 7): LEARNINGS SO FAR
    # ==========================================
    inject_content(prs.slides[7], "LEARNINGS", [
        ("New Skills Learned:", 0),
        ("Integrating external Large Language Models (Gemini) into full-stack web applications.", 1),
        ("Managing state and session cookies securely via a Flask backend.", 1),
        ("Advanced Prompt Engineering to bypass AI 'hallucinations' and enforce data structures.", 1),
        ("Solution Demo URL:", 0),
        ("Localhost: http://127.0.0.1:5000", 1),
        ("[INSERT IMAGE: PASTE A SCREENSHOT OF THE APP RUNNING HERE]", 0)
    ])

    output_filename = 'AI_Quiz_Arena_Final_Submission.pptx'
    prs.save(output_filename)
    print(f"Success! Your presentation is fully packed with info. Open '{output_filename}' to review.")

if __name__ == '__main__':
    fill_university_template()